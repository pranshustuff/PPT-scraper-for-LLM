#!/usr/bin/env python3

import os
import json
import requests
from pptx import Presentation
from itertools import islice

LOCAL_PPT_DIR = "ppt_files"
os.makedirs(LOCAL_PPT_DIR, exist_ok=True)

OUTPUT_DIR = "ppt_summaries"
os.makedirs(OUTPUT_DIR, exist_ok=True)

CHUNK_SIZE = 5

ppt_urls = [
    # "https://www.example.com/sample1.pptx",
]

def download_ppt(url, save_dir=LOCAL_PPT_DIR):
    local_filename = os.path.join(save_dir, url.split("/")[-1])
    try:
        r = requests.get(url)
        r.raise_for_status()
        with open(local_filename, "wb") as f:
            f.write(r.content)
        print(f"[INFO] Downloaded {local_filename}")
        return local_filename
    except Exception as e:
        print(f"[WARN] Failed to download {url}: {e}")
        return None

def scan_local_folder(folder=LOCAL_PPT_DIR):
    ppt_files = [os.path.join(folder, f) for f in os.listdir(folder)
                 if f.lower().endswith(".pptx")]
    print(f"[INFO] Found {len(ppt_files)} PPTX files locally")
    return ppt_files

def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    slides_text = []
    for slide in prs.slides:
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())
        slides_text.append("\n".join(slide_content))
    return slides_text

def chunked(iterable, size=CHUNK_SIZE):
    it = iter(iterable)
    while True:
        chunk = list(islice(it, size))
        if not chunk:
            break
        yield chunk

def summarize_text(text, max_tokens=250):
    prompt = f"Summarize the following lecture slides concisely:\n\n{text}"
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=max_tokens,
        temperature=0.3
    )
    return response.choices[0].message["content"]

def process_ppt(ppt_path):
    slides_text = extract_text_from_ppt(ppt_path)
    summaries = []
    for i, chunk in enumerate(chunked(slides_text)):
        chunk_text = "\n\n".join(chunk)
        summary = summarize_text(chunk_text)
        summaries.append({
            "chunk_index": i,
            "slides": chunk,
            "summary": summary
        })
    return summaries

def save_summaries_to_json(ppt_path, summaries, output_dir=OUTPUT_DIR):
    base_name = os.path.splitext(os.path.basename(ppt_path))[0]
    output_file = os.path.join(output_dir, f"{base_name}_summary.json")
    with open(output_file, "w", encoding="utf-8") as f:
        json.dump(summaries, f, indent=4, ensure_ascii=False)
    print(f"[INFO] Saved summary to {output_file}")

if __name__ == "__main__":
    
    for url in ppt_urls:
        download_ppt(url)
    ppt_files = scan_local_folder()
    for ppt in ppt_files:
        print(f"[INFO] Processing {ppt}")
        summaries = process_ppt(ppt)
        save_summaries_to_json(ppt, summaries)

    print("[INFO] All done! You can now feed JSON summaries to NotebookLM or other LLM pipelines.")
